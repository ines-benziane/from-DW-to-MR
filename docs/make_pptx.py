"""Script de génération de la présentation PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

BLUE   = RGBColor(0x00, 0x56, 0xB3)
LBLUE  = RGBColor(0xE8, 0xF0, 0xFB)
DGRAY  = RGBColor(0x33, 0x33, 0x33)
GRAY   = RGBColor(0x88, 0x88, 0x88)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def header(slide, title, subtitle=None):
    """Bande bleue en haut avec titre."""
    add_rect(slide, 0, 0, W, Inches(1.1), fill=BLUE)
    add_text(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.7),
             size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.78), Inches(12.5), Inches(0.35),
                 size=13, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)


def placeholder_box(slide, l, t, w, h, label="[ Capture d'écran ]"):
    add_rect(slide, l, t, w, h, fill=LBLUE, line=BLUE)
    add_text(slide, label, l, t + h//2 - Inches(0.2), w, Inches(0.4),
             size=11, color=GRAY, align=PP_ALIGN.CENTER)


def footer(slide, text="Institut de Myologie — IRM musculaire"):
    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), fill=RGBColor(0xF0, 0xF4, 0xFF))
    add_text(slide, text, Inches(0.3), H - Inches(0.32), Inches(12.7), Inches(0.3),
             size=9, color=GRAY, align=PP_ALIGN.LEFT)


# ── helpers ───────────────────────────────────────────────────────────────────

def bullet(slide, items, l, t, w, size=14, spacing=0.38):
    for i, (txt, sub) in enumerate(items):
        y = t + Inches(spacing * i)
        add_text(slide, "▸  " + txt, l, y, w, Inches(0.35),
                 size=size, bold=True, color=BLUE)
        if sub:
            add_text(slide, "     " + sub, l, y + Inches(0.22), w, Inches(0.28),
                     size=11, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, fill=BLUE)
    add_rect(sl, Inches(0.5), Inches(2.0), Inches(12.3), Inches(3.5),
             fill=WHITE)
    add_text(sl, "Comptes-rendus IRM musculaires",
             Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.0),
             size=32, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(sl, "Présentation des versions, palettes de couleur et commentaire automatique",
             Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.6),
             size=16, color=DGRAY, align=PP_ALIGN.CENTER)
    add_text(sl, "Institut de Myologie — Juin 2026",
             Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.4),
             size=12, color=GRAY, align=PP_ALIGN.CENTER)


def slide_versions_tableau(prs):
    sl = blank_slide(prs)
    header(sl, "Les 6 versions de compte-rendu")
    footer(sl)

    rows = [
        ("1A", "1 coupe volumétrique",      "Non",     "Vue transversale uniquement"),
        ("1B", "1 coupe volumétrique",      "Non",     "Vues de face + transversale (cuisses)"),
        ("2A", "5 coupes — format compact", "Synth. v1", "Couleurs FF sans labels"),
        ("2B", "5 coupes — format compact", "Synth. v2", "Couleurs FF sans labels"),
        ("3A", "5 coupes — format complet", "Synth. v1", "FF avec labels des muscles"),
        ("3B", "5 coupes — format complet", "Synth. v2", "FF avec labels des muscles"),
    ]
    cols = ["Version", "Coupes", "Synthèse", "FF cuisses"]
    col_w = [Inches(1.1), Inches(3.5), Inches(2.0), Inches(5.5)]
    col_x = [Inches(0.4), Inches(1.55), Inches(5.1), Inches(7.15)]
    row_h = Inches(0.52)
    top   = Inches(1.3)

    # header row
    for j, (cx, cw, col) in enumerate(zip(col_x, col_w, cols)):
        add_rect(sl, cx, top, cw - Inches(0.05), row_h, fill=BLUE)
        add_text(sl, col, cx + Inches(0.05), top + Inches(0.1), cw, row_h,
                 size=13, bold=True, color=WHITE)

    for i, (ver, coupes, synth, ff) in enumerate(rows):
        y = top + row_h + Inches(0.52 * i)
        bg = LBLUE if i % 2 == 0 else WHITE
        vals = [ver, coupes, synth, ff]
        for j, (cx, cw, val) in enumerate(zip(col_x, col_w, vals)):
            add_rect(sl, cx, y, cw - Inches(0.05), row_h, fill=bg, line=RGBColor(0xCC,0xDD,0xFF))
            bold = (j == 0)
            col_c = BLUE if j == 0 else DGRAY
            add_text(sl, val, cx + Inches(0.08), y + Inches(0.1), cw, row_h,
                     size=12, bold=bold, color=col_c)

    add_text(sl, "* La version A et B diffèrent uniquement par la mise en page de la page de synthèse.",
             Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3),
             size=10, color=GRAY)


def slide_version_1(prs):
    sl = blank_slide(prs)
    header(sl, "Versions 1A / 1B — 1 coupe volumétrique", "Sans page de synthèse dédiée")
    footer(sl)

    items = [
        ("1 coupe par segment", "La valeur affichée est la moyenne sur tout le volume musculaire."),
        ("Comparaison avant / après", "Si un examen antérieur est disponible : coupe antécédente + évolution par muscle."),
        ("Commentaire T2 automatique", "Catégorie générée selon le nombre de muscles ≥ 39 ms."),
        ("1B uniquement — vues de face FF cuisses", "Vues anatomiques antérieure et postérieure des cuisses."),
    ]
    bullet(sl, items, Inches(0.5), Inches(1.3), Inches(5.8), size=13, spacing=0.85)
    placeholder_box(sl, Inches(6.8), Inches(1.2), Inches(6.1), Inches(5.8),
                    "[ Capture 1A ou 1B ]")


def slide_versions_25(prs):
    sl = blank_slide(prs)
    header(sl, "Versions 2 et 3 — 5 coupes", "Format compact (2) vs complet (3)")
    footer(sl)

    add_text(sl, "Format compact — 2A / 2B", Inches(0.4), Inches(1.25), Inches(6.0), Inches(0.4),
             size=15, bold=True, color=BLUE)
    items_2 = [
        ("5 coupes (1 sur 2)", "Présentation condensée, idéale pour une lecture rapide."),
        ("FF : couleurs uniquement", "Pas de labels de muscles affichés."),
    ]
    bullet(sl, items_2, Inches(0.5), Inches(1.7), Inches(5.8), size=12, spacing=0.72)

    add_text(sl, "Format complet — 3A / 3B", Inches(0.4), Inches(3.3), Inches(6.0), Inches(0.4),
             size=15, bold=True, color=BLUE)
    items_3 = [
        ("5 coupes (1 sur 2)", "Même sélection de coupes que le format compact."),
        ("FF : avec labels des muscles", "Noms des muscles visibles sur chaque coupe."),
    ]
    bullet(sl, items_3, Inches(0.5), Inches(3.75), Inches(5.8), size=12, spacing=0.72)

    placeholder_box(sl, Inches(6.8), Inches(1.2), Inches(6.1), Inches(5.8),
                    "[ Capture 2A ou 3A ]")


def slide_synthese(prs):
    sl = blank_slide(prs)
    header(sl, "Page de synthèse — v1 et v2", "Générée uniquement si un examen antérieur est disponible")
    footer(sl)

    add_text(sl, "Synthèse v1", Inches(0.4), Inches(1.25), Inches(5.5), Inches(0.4),
             size=15, bold=True, color=BLUE)
    add_text(sl, "Une section par biomarqueur / segment, avec évolution par muscle.",
             Inches(0.5), Inches(1.65), Inches(5.5), Inches(0.5), size=12, color=DGRAY)

    add_text(sl, "Synthèse v2", Inches(0.4), Inches(2.5), Inches(5.5), Inches(0.4),
             size=15, bold=True, color=BLUE)
    add_text(sl, "Blocs T2 / FF séparés : examen actuel puis évolution distinctement présentés.",
             Inches(0.5), Inches(2.9), Inches(5.5), Inches(0.5), size=12, color=DGRAY)

    add_text(sl, "Dans les deux cas :", Inches(0.4), Inches(3.75), Inches(5.5), Inches(0.35),
             size=13, bold=True, color=DGRAY)
    items = [
        ("T2 : variation en ms", "(ex. : +6 ms)"),
        ("FF : variation en points de %", "(ex. : +7 %)"),
        ("Couleur d'évolution", "Bleu = amélioration · Rouge = aggravation · Seuil ±5 unités"),
    ]
    bullet(sl, items, Inches(0.5), Inches(4.15), Inches(5.8), size=12, spacing=0.72)

    placeholder_box(sl, Inches(6.8), Inches(1.2), Inches(6.1), Inches(5.8),
                    "[ Capture page de synthèse ]")


def slide_colormaps(prs):
    sl = blank_slide(prs)
    header(sl, "Palettes de couleur")
    footer(sl)

    add_text(sl, "T2 — 4 palettes disponibles",
             Inches(0.4), Inches(1.25), Inches(12.0), Inches(0.4),
             size=15, bold=True, color=BLUE)

    t2_items = [
        ("default (vikO + lajolla)", "Bleu → jaune à 39 ms → rouge. Référence actuelle."),
        ("hawaii_r",                 "Dégradé chaud, même logique jaune autour du seuil."),
        ("roma_r",                   "Du bleu foncé au rouge vif."),
        ("bam_r",                    "Palette divergente froide → chaude."),
    ]
    bullet(sl, t2_items, Inches(0.5), Inches(1.7), Inches(12.0), size=12, spacing=0.65)

    add_text(sl, "FF — 3 palettes disponibles",
             Inches(0.4), Inches(4.5), Inches(12.0), Inches(0.4),
             size=15, bold=True, color=BLUE)
    ff_items = [
        ("default (lajolla_r)", "Référence actuelle — coupée à 60 % de la palette pour les muscles très infiltrés."),
        ("lapaz",               "Palette froide à chaude sur toute la plage."),
        ("davos",               "Palette gris clair à sombre."),
    ]
    bullet(sl, ff_items, Inches(0.5), Inches(4.95), Inches(12.0), size=12, spacing=0.65)


def slide_commentaire(prs):
    sl = blank_slide(prs)
    header(sl, "Commentaire automatique T2")
    footer(sl)

    add_text(sl,
             "Un commentaire est proposé à la fin de chaque section T2. "
             "Il est basé sur le nombre de muscles dont le T2 volumétrique dépasse 39 ms.",
             Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.55),
             size=13, color=DGRAY)

    cats = [
        ("Normal strict",  "Aucune coupe ≥ 39 ms, valeur max nettement en dessous du seuil.",
         "« Valeurs des moyennes du T2 musculaire dans la norme. »"),
        ("Normal limite",  "Aucun muscle ≥ 39 ms, mais certaines coupes atteignent 39 ms.",
         "« Valeurs de T2 dans la norme, suggérant que les processus inflammatoires et lésionnels sont peu actifs. »"),
        ("Discret",        "1 muscle ≥ 39 ms.",
         "« Discrète augmentation de la moyenne du T2 musculaire en regard de [muscle]… »"),
        ("Modéré",         "2 à 3 muscles ≥ 39 ms.",
         "« Augmentation modérée et significative des moyennes du T2 musculaire… »"),
        ("Significatif",   "4 muscles ou plus ≥ 39 ms.",
         "« Augmentation significative des moyennes du T2 musculaire… »"),
    ]

    row_h = Inches(0.88)
    top   = Inches(1.85)
    for i, (cat, cond, comment) in enumerate(cats):
        y  = top + Inches(0.9 * i)
        bg = LBLUE if i % 2 == 0 else WHITE
        add_rect(sl, Inches(0.3), y, Inches(12.7), row_h - Inches(0.05),
                 fill=bg, line=RGBColor(0xCC, 0xDD, 0xFF))
        add_text(sl, cat, Inches(0.4), y + Inches(0.05), Inches(2.0), Inches(0.35),
                 size=12, bold=True, color=BLUE)
        add_text(sl, cond, Inches(2.45), y + Inches(0.05), Inches(4.5), Inches(0.35),
                 size=11, color=DGRAY)
        add_text(sl, comment, Inches(7.0), y + Inches(0.05), Inches(5.8), Inches(0.75),
                 size=10, color=GRAY)


def slide_conclusion(prs):
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, fill=BLUE)
    add_rect(sl, Inches(0.5), Inches(2.2), Inches(12.3), Inches(3.0), fill=WHITE)

    add_text(sl, "Merci pour votre attention",
             Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.8),
             size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(sl,
             "Ces versions sont disponibles pour ceux qui souhaitent y jeter un œil.\n"
             "Vos retours guideront les prochaines évolutions.",
             Inches(0.8), Inches(3.35), Inches(11.7), Inches(0.8),
             size=14, color=DGRAY, align=PP_ALIGN.CENTER)
    add_text(sl, "Inès Benziane — Institut de Myologie",
             Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.4),
             size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    slide_title(prs)
    slide_versions_tableau(prs)
    slide_version_1(prs)
    slide_versions_25(prs)
    slide_synthese(prs)
    slide_colormaps(prs)
    slide_commentaire(prs)
    slide_conclusion(prs)
    out = "docs/presentation_CR_IRM.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
