"""Présentation école d'été — IRM musculaire, compte-rendu standardisé."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE   = RGBColor(0x00, 0x56, 0xB3)
DBLUE  = RGBColor(0x00, 0x3A, 0x7A)
LBLUE  = RGBColor(0xE8, 0xF0, 0xFB)
DGRAY  = RGBColor(0x33, 0x33, 0x33)
GRAY   = RGBColor(0x88, 0x88, 0x88)
LGRAY  = RGBColor(0xF4, 0xF4, 0xF4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xF5, 0xC4, 0x00)
RED    = RGBColor(0xCC, 0x33, 0x00)
GREEN  = RGBColor(0x27, 0xAE, 0x60)

W = Inches(13.33)
H = Inches(7.5)


# ── primitives ─────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(0.75)):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size=14, bold=False, italic=False,
        color=RGBColor(0x33,0x33,0x33), align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(1.1), fill=BLUE)
    txt(slide, title, Inches(0.4), Inches(0.13), Inches(12.5), Inches(0.72),
        size=24, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.4), Inches(0.78), Inches(12.5), Inches(0.35),
            size=12, italic=True, color=RGBColor(0xBB, 0xCC, 0xEE))

def footer(slide, note=None):
    rect(slide, 0, H - Inches(0.35), W, Inches(0.35), fill=LGRAY)
    label = note or "Institut de Myologie — Hôpital La Pitié-Salpêtrière"
    txt(slide, label, Inches(0.3), H - Inches(0.32), Inches(12.7), Inches(0.3),
        size=9, color=GRAY)

def placeholder(slide, l, t, w, h, label="[ Capture d'écran ]"):
    rect(slide, l, t, w, h, fill=LBLUE, line_color=BLUE)
    txt(slide, label, l, t + h//2 - Inches(0.2), w, Inches(0.4),
        size=11, color=GRAY, align=PP_ALIGN.CENTER)

def bullet_list(slide, items, l, t, w, size=13, gap=0.55):
    """items = list of (title, subtitle) or (title, None)"""
    for i, (title, sub) in enumerate(items):
        y = t + Inches(gap * i)
        txt(slide, "▸  " + title, l, y, w, Inches(0.35),
            size=size, bold=True, color=BLUE)
        if sub:
            txt(slide, "     " + sub, l, y + Inches(0.24), w, Inches(0.28),
                size=size - 2, color=DGRAY)

def tag(slide, label, l, t, w=Inches(2.2), h=Inches(0.38), bg=BLUE, fg=WHITE):
    rect(slide, l, t, w, h, fill=bg)
    txt(slide, label, l + Inches(0.1), t + Inches(0.04), w, h,
        size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def s_titre(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, fill=DBLUE)
    rect(sl, 0, Inches(2.0), W, Inches(3.8), fill=WHITE)

    txt(sl, "Compte-rendu IRM musculaire standardisé",
        Inches(0.7), Inches(2.2), Inches(11.9), Inches(1.0),
        size=30, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    txt(sl, "Conception, choix de représentation et automatisation",
        Inches(0.7), Inches(3.3), Inches(11.9), Inches(0.55),
        size=16, color=DGRAY, align=PP_ALIGN.CENTER)

    txt(sl, "École d'été — Neurologie 2026",
        Inches(0.7), Inches(4.05), Inches(11.9), Inches(0.4),
        size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    txt(sl, "Inès Benziane — Institut de Myologie, Paris",
        Inches(0.7), Inches(5.3), Inches(11.9), Inches(0.4),
        size=12, color=WHITE, align=PP_ALIGN.CENTER)


def s_contexte(prs):
    sl = blank_slide(prs)
    header(sl, "Contexte clinique", "Pourquoi un CR standardisé pour l'IRM musculaire ?")
    footer(sl)

    items = [
        ("Un besoin de standardisation",
         "Les maladies neuromusculaires nécessitent un suivi longitudinal rigoureux "
         "des muscles par IRM. Sans CR structuré, l'interprétation dépend fortement de l'opérateur."),
        ("Deux biomarqueurs complémentaires",
         "Le temps de relaxation T2 (inflammation, lésions actives) "
         "et la fraction de graisse FF (infiltration lipidique irréversible)."),
        ("Objectif",
         "Proposer un CR lisible, reproductible, comparable dans le temps — "
         "et qui s'adapte aux différents contextes de lecture."),
    ]

    for i, (title, body) in enumerate(items):
        y = Inches(1.3 + 1.65 * i)
        rect(sl, Inches(0.4), y, Inches(12.5), Inches(1.45), fill=LBLUE, line_color=BLUE)
        txt(sl, title, Inches(0.6), y + Inches(0.1), Inches(12.0), Inches(0.4),
            size=13, bold=True, color=BLUE)
        txt(sl, body, Inches(0.6), y + Inches(0.48), Inches(12.0), Inches(0.85),
            size=12, color=DGRAY)


def s_t2(prs):
    sl = blank_slide(prs)
    header(sl, "Le biomarqueur T2", "Temps de relaxation transversal — marqueur d'inflammation et de lésion active")
    footer(sl)

    txt(sl, "Le T2 musculaire est prolongé en présence d'un œdème ou d'une inflammation active. "
           "Il constitue un indicateur sensible des processus lésionnels en cours.",
        Inches(0.4), Inches(1.2), Inches(7.5), Inches(0.7), size=12, color=DGRAY)

    txt(sl, "Échelle de couleur retenue", Inches(0.4), Inches(2.0), Inches(7.5), Inches(0.35),
        size=13, bold=True, color=BLUE)

    scale = [
        (GRAY,   "< 10 ms",    "Incertitude de mesure — non interprétable"),
        (BLUE,   "10 – 37 ms", "Plage normale"),
        (YELLOW, "37 – 41 ms", "Zone de vigilance — seuil à 39 ms"),
        (RED,    "> 41 ms",    "Valeur élevée — atteinte probable"),
    ]
    for i, (color, val, label) in enumerate(scale):
        y = Inches(2.45 + 0.78 * i)
        rect(sl, Inches(0.4), y, Inches(0.35), Inches(0.55), fill=color)
        txt(sl, val,   Inches(0.85), y + Inches(0.08), Inches(1.8),  Inches(0.38), size=11, bold=True,  color=DGRAY)
        txt(sl, label, Inches(2.7),  y + Inches(0.08), Inches(5.5),  Inches(0.38), size=11, color=DGRAY)

    txt(sl, "Pourquoi 39 ms ?",
        Inches(0.4), Inches(5.65), Inches(7.5), Inches(0.35),
        size=12, bold=True, color=BLUE)
    txt(sl, "Seuil issu de la littérature et validé en interne — "
           "les muscles dont le T2 dépasse ce seuil sont mis en évidence par un contour plus épais.",
        Inches(0.4), Inches(6.0), Inches(7.5), Inches(0.55), size=11, color=DGRAY)

    placeholder(sl, Inches(8.1), Inches(1.2), Inches(4.8), Inches(5.8),
                "[ Capture section T2 ]")


def s_ff(prs):
    sl = blank_slide(prs)
    header(sl, "La fraction de graisse (FF)", "Marqueur d'infiltration lipidique — processus irréversible")
    footer(sl)

    txt(sl, "La FF mesure la proportion de graisse dans le muscle. "
           "Contrairement au T2, elle reflète des lésions établies et non récupérables au-delà d'un certain seuil.",
        Inches(0.4), Inches(1.2), Inches(7.5), Inches(0.65), size=12, color=DGRAY)

    txt(sl, "Trois zones cliniques", Inches(0.4), Inches(2.0), Inches(7.5), Inches(0.35),
        size=13, bold=True, color=BLUE)

    zones = [
        (GREEN, "0 – 5 %",  "Infiltration très faible — couleur uniforme."),
        (BLUE,  "5 – 40 %", "Zone d'intérêt clinique — nuances de couleur pour distinguer les degrés."),
        (RED,   "> 40 %",   "Infiltration sévère — muscle non récupérable. Couleur uniforme saturée."),
    ]
    for i, (color, val, label) in enumerate(zones):
        y = Inches(2.45 + 0.95 * i)
        rect(sl, Inches(0.4), y, Inches(0.35), Inches(0.7), fill=color)
        txt(sl, val,   Inches(0.85), y + Inches(0.1), Inches(1.8), Inches(0.38), size=12, bold=True, color=DGRAY)
        txt(sl, label, Inches(2.7),  y + Inches(0.1), Inches(5.2), Inches(0.55), size=11, color=DGRAY)

    txt(sl, "La FF est calculée par la méthode Dixon 3 points. "
           "Elle est exprimée en proportion (0 à 100 %).",
        Inches(0.4), Inches(5.6), Inches(7.5), Inches(0.55), size=11, italic=True, color=GRAY)

    placeholder(sl, Inches(8.1), Inches(1.2), Inches(4.8), Inches(5.8),
                "[ Capture section FF ]")


def s_structure(prs):
    sl = blank_slide(prs)
    header(sl, "Structure du compte-rendu")
    footer(sl)

    sections = [
        ("En-tête",              "Informations patient, date d'examen, équipe médicale."),
        ("Repères anatomiques",  "Schéma de coupe avec légende des abréviations musculaires."),
        ("Section T2",           "Coupes colorées + colorbar + commentaire automatique."),
        ("Section FF",           "Coupes colorées + vues anatomiques de face (cuisses)."),
        ("Page de synthèse",     "Évolution par muscle entre deux examens (si antécédent disponible)."),
    ]

    for i, (title, desc) in enumerate(sections):
        y  = Inches(1.25 + 1.02 * i)
        bg = LBLUE if i % 2 == 0 else WHITE
        rect(sl, Inches(0.4), y, Inches(7.5), Inches(0.9), fill=bg, line_color=BLUE)
        tag(sl, f"0{i+1}", Inches(0.45), y + Inches(0.17), w=Inches(0.5), h=Inches(0.55), bg=BLUE)
        txt(sl, title, Inches(1.1), y + Inches(0.08), Inches(3.5), Inches(0.38),
            size=13, bold=True, color=BLUE)
        txt(sl, desc,  Inches(1.1), y + Inches(0.46), Inches(6.6), Inches(0.38),
            size=11, color=DGRAY)

    placeholder(sl, Inches(8.1), Inches(1.2), Inches(4.8), Inches(5.8),
                "[ Vue d'ensemble du CR ]")


def s_segmentation(prs):
    sl = blank_slide(prs)
    header(sl, "Segmentation automatique des muscles",
           "Les contours musculaires sont générés automatiquement par un algorithme dédié")
    footer(sl)

    txt(sl, "Chaque muscle est délimité automatiquement sur chaque coupe IRM. "
           "Ces contours sont ensuite lissés (interpolation B-spline) avant affichage dans le CR.",
        Inches(0.4), Inches(1.2), Inches(7.5), Inches(0.65), size=12, color=DGRAY)

    items = [
        ("Ce que ça apporte",
         "Reproductibilité totale entre patients et entre examens. "
         "Gain de temps considérable par rapport à une segmentation manuelle."),
        ("Limites à connaître",
         "La qualité des contours dépend de la qualité de la segmentation automatique. "
         "Une erreur de segmentation se répercute directement sur les valeurs T2 et FF affichées."),
        ("Mention dans le CR",
         "Une note indique explicitement que les contours sont issus d'une segmentation automatique."),
    ]

    for i, (title, body) in enumerate(items):
        y = Inches(2.0 + 1.45 * i)
        rect(sl, Inches(0.4), y, Inches(7.5), Inches(1.3), fill=LBLUE, line_color=BLUE)
        txt(sl, title, Inches(0.6), y + Inches(0.08), Inches(7.0), Inches(0.38),
            size=13, bold=True, color=BLUE)
        txt(sl, body,  Inches(0.6), y + Inches(0.44), Inches(7.0), Inches(0.75),
            size=11, color=DGRAY)

    placeholder(sl, Inches(8.1), Inches(1.2), Inches(4.8), Inches(5.8),
                "[ Capture contours musculaires ]")


def s_colormap_t2(prs):
    sl = blank_slide(prs)
    header(sl, "Choix de la colormap T2", "Rendre visible ce qui est cliniquement pertinent")
    footer(sl)

    blocs = [
        ("Gris en dessous de 10 ms",
         "Valeurs trop basses : incertitude de mesure connue à ces niveaux. "
         "Le gris signale visuellement l'absence d'interprétation fiable."),
        ("Bleu → bleu clair (10 – 37 ms)",
         "Plage normale. Le dégradé permet de voir les variations internes "
         "sans alarmer le lecteur."),
        ("Jaune autour de 39 ms (37 – 41 ms)",
         "Zone de vigilance. Le jaune attire l'œil sur les muscles proches du seuil "
         "sans être aussi alarmant que le rouge."),
        ("Orange → rouge au-delà de 41 ms",
         "Signal d'alerte fort. De plus, le contour du muscle est épaissi "
         "pour renforcer la mise en évidence."),
    ]
    for i, (title, body) in enumerate(blocs):
        y = Inches(1.25 + 1.4 * i)
        colors = [GRAY, BLUE, YELLOW, RED]
        rect(sl, Inches(0.4),  y, Inches(0.15), Inches(1.2), fill=colors[i])
        rect(sl, Inches(0.6),  y, Inches(7.3),  Inches(1.2), fill=LBLUE if i%2==0 else WHITE,
             line_color=RGBColor(0xCC,0xDD,0xFF))
        txt(sl, title, Inches(0.75), y + Inches(0.08), Inches(7.0), Inches(0.38),
            size=12, bold=True, color=BLUE)
        txt(sl, body,  Inches(0.75), y + Inches(0.44), Inches(7.0), Inches(0.65),
            size=11, color=DGRAY)

    placeholder(sl, Inches(8.1), Inches(1.2), Inches(4.8), Inches(5.8),
                "[ Colorbar T2 + exemple ]")


def s_colormap_ff(prs):
    sl = blank_slide(prs)
    header(sl, "Choix de la colormap FF", "Concentrer l'attention sur la zone d'intérêt clinique")
    footer(sl)

    txt(sl, "La palette FF est construite autour de 3 zones aux logiques distinctes :",
        Inches(0.4), Inches(1.2), Inches(7.5), Inches(0.4), size=13, bold=True, color=DGRAY)

    blocs = [
        (GREEN, "0 – 5 % — zone basse",
         "Infiltration très faible. Une couleur uniforme suffit : il n'y a pas de nuance "
         "cliniquement significative à cette échelle."),
        (BLUE, "5 – 40 % — zone d'intérêt",
         "C'est ici que les nuances comptent. "
         "La palette est étalée sur cette plage pour maximiser la lisibilité des différences."),
        (RED, "> 40 % — zone sévère",
         "Au-delà de ce seuil, le muscle est considéré comme non récupérable. "
         "Une couleur uniforme suffit — la nuance n'apporte plus d'information utile."),
    ]
    for i, (color, title, body) in enumerate(blocs):
        y = Inches(1.75 + 1.5 * i)
        rect(sl, Inches(0.4), y, Inches(0.15), Inches(1.3), fill=color)
        rect(sl, Inches(0.6), y, Inches(7.3),  Inches(1.3),
             fill=LBLUE if i % 2 == 0 else WHITE,
             line_color=RGBColor(0xCC, 0xDD, 0xFF))
        txt(sl, title, Inches(0.75), y + Inches(0.08), Inches(7.0), Inches(0.38),
            size=12, bold=True, color=BLUE)
        txt(sl, body,  Inches(0.75), y + Inches(0.46), Inches(7.0), Inches(0.72),
            size=11, color=DGRAY)

    placeholder(sl, Inches(8.1), Inches(1.2), Inches(4.8), Inches(5.8),
                "[ Colorbar FF + exemple ]")


def s_commentaire(prs):
    sl = blank_slide(prs)
    header(sl, "Commentaire automatique T2",
           "Une aide à la lecture — non un diagnostic")
    footer(sl)

    txt(sl, "Un commentaire textuel est généré automatiquement à la fin de chaque section T2, "
           "selon le nombre de muscles dont le T2 volumétrique dépasse 39 ms.",
        Inches(0.4), Inches(1.2), Inches(8.5), Inches(0.55), size=12, color=DGRAY)

    cats = [
        ("Normal strict",  "Aucune coupe ≥ 39 ms"),
        ("Normal limite",  "Aucun muscle ≥ 39 ms, mais coupes proches du seuil"),
        ("Discret",        "1 muscle ≥ 39 ms"),
        ("Modéré",         "2 à 3 muscles ≥ 39 ms"),
        ("Significatif",   "4 muscles ou plus ≥ 39 ms"),
    ]
    for i, (cat, cond) in enumerate(cats):
        y  = Inches(1.85 + 0.82 * i)
        bg = LBLUE if i % 2 == 0 else WHITE
        rect(sl, Inches(0.4), y, Inches(8.5), Inches(0.72),
             fill=bg, line_color=RGBColor(0xCC, 0xDD, 0xFF))
        txt(sl, cat,  Inches(0.55), y + Inches(0.12), Inches(2.2), Inches(0.45),
            size=12, bold=True, color=BLUE)
        txt(sl, cond, Inches(2.8),  y + Inches(0.12), Inches(6.0), Inches(0.45),
            size=12, color=DGRAY)

    rect(sl, Inches(0.4), Inches(6.2), Inches(8.5), Inches(0.75),
         fill=RGBColor(0xFF, 0xF8, 0xE6), line_color=RGBColor(0xF0, 0xC0, 0x00))
    txt(sl, "⚠  Ce commentaire est une aide à la lecture. Il ne tient pas compte du contexte "
           "clinique, de l'histoire du patient ni des valeurs FF. Le médecin reste décisionnaire.",
        Inches(0.55), Inches(6.25), Inches(8.2), Inches(0.65), size=11, italic=True, color=DGRAY)


def s_synthese(prs):
    sl = blank_slide(prs)
    header(sl, "Suivi longitudinal — page de synthèse",
           "Générée uniquement si un examen antérieur est disponible")
    footer(sl)

    txt(sl, "La page de synthèse présente l'évolution de chaque muscle entre deux examens. "
           "Elle permet un suivi quantitatif et visuel de la progression de la maladie.",
        Inches(0.4), Inches(1.2), Inches(7.5), Inches(0.6), size=12, color=DGRAY)

    items = [
        ("T2 : variation en ms",
         "Différence absolue entre les deux examens (ex. : +6 ms)."),
        ("FF : variation en points de %",
         "Différence absolue (ex. : passage de 20 % à 27 % = +7 %)."),
        ("Couleur d'évolution",
         "Bleu = amélioration · Blanc = stable (±5 unités) · Rouge = aggravation."),
        ("Représentation visuelle",
         "Chaque muscle est coloré selon son sens d'évolution sur la coupe anatomique."),
    ]
    bullet_list(sl, items, Inches(0.4), Inches(1.95), Inches(7.5), size=12, gap=0.88)

    placeholder(sl, Inches(8.1), Inches(1.2), Inches(4.8), Inches(5.8),
                "[ Capture page de synthèse ]")


def s_perspectives(prs):
    sl = blank_slide(prs)
    header(sl, "Perspectives")
    footer(sl)

    items = [
        ("Validation clinique des seuils",
         "Les seuils T2 (39 ms) et FF (40 %) sont à valider sur une cohorte plus large."),
        ("Amélioration du commentaire automatique",
         "Intégrer les valeurs FF et le contexte longitudinal dans la génération du commentaire."),
        ("Extension à d'autres segments",
         "Le pipeline est actuellement validé sur cuisses et jambes."),
        ("Retours des utilisateurs",
         "Ajustements en cours suite aux premières lectures par l'équipe médicale."),
    ]
    bullet_list(sl, items, Inches(0.8), Inches(1.4), Inches(11.7), size=14, gap=1.2)


def s_conclusion(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, fill=DBLUE)
    rect(sl, Inches(0.5), Inches(1.8), Inches(12.3), Inches(4.0), fill=WHITE)

    txt(sl, "En résumé",
        Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.6),
        size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    points = [
        "Un CR standardisé, lisible et comparable dans le temps.",
        "Deux biomarqueurs : T2 (lésion active) + FF (infiltration irréversible).",
        "Segmentation automatique + colormaps + commentaire automatique.",
        "Un outil évolutif — vos retours sont précieux.",
    ]
    for i, p in enumerate(points):
        txt(sl, "✓  " + p,
            Inches(1.2), Inches(2.65 + 0.56 * i), Inches(11.0), Inches(0.45),
            size=13, color=DGRAY, align=PP_ALIGN.LEFT)

    txt(sl, "Merci pour votre attention — questions bienvenues",
        Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
        size=14, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "Inès Benziane — Institut de Myologie, Paris",
        Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
        size=11, color=RGBColor(0xBB, 0xCC, 0xEE), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    s_titre(prs)
    s_contexte(prs)
    s_t2(prs)
    s_ff(prs)
    s_structure(prs)
    s_segmentation(prs)
    s_colormap_t2(prs)
    s_colormap_ff(prs)
    s_commentaire(prs)
    s_synthese(prs)
    s_perspectives(prs)
    s_conclusion(prs)

    out = "docs/presentation_ecole_ete.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
